"""
Multi-format document parser.

Supports: PDF, DOCX, PPTX, XLSX/CSV, TXT, Markdown, HTML, JSON
Extracts clean text + metadata from any supported format.
"""

import os
import json
import chardet
import logging
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional

logger = logging.getLogger(__name__)


@dataclass
class ParsedDocument:
    """Represents a parsed document with its text content and metadata."""
    content: str
    metadata: dict = field(default_factory=dict)
    filename: str = ""
    file_type: str = ""
    page_count: int = 0


class DocumentParser:
    """
    Universal document parser that handles multiple file formats.
    Each format has a dedicated parser for maximum extraction quality.
    """

    SUPPORTED_EXTENSIONS = {
        ".pdf", ".docx", ".doc", ".pptx", ".xlsx", ".xls",
        ".csv", ".txt", ".md", ".html", ".htm", ".json",
        ".rtf", ".log", ".xml"
    }

    def parse(self, file_path: str) -> ParsedDocument:
        """
        Parse a document from the given file path.
        Automatically detects the format and uses the appropriate parser.
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = path.suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file format: {ext}. "
                f"Supported: {', '.join(sorted(self.SUPPORTED_EXTENSIONS))}"
            )

        logger.info(f"Parsing {path.name} (format: {ext})")

        parser_map = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".doc": self._parse_docx,
            ".pptx": self._parse_pptx,
            ".xlsx": self._parse_excel,
            ".xls": self._parse_excel,
            ".csv": self._parse_csv,
            ".txt": self._parse_text,
            ".log": self._parse_text,
            ".md": self._parse_markdown,
            ".html": self._parse_html,
            ".htm": self._parse_html,
            ".json": self._parse_json,
            ".xml": self._parse_html,
            ".rtf": self._parse_text,
        }

        parser_fn = parser_map.get(ext, self._parse_text)
        doc = parser_fn(file_path)
        doc.filename = path.name
        doc.file_type = ext
        doc.metadata["source"] = path.name
        doc.metadata["file_type"] = ext
        doc.metadata["file_size_bytes"] = path.stat().st_size

        # Clean up extracted text
        doc.content = self._clean_text(doc.content)

        logger.info(
            f"Parsed {path.name}: {len(doc.content)} chars, "
            f"{doc.page_count} pages"
        )
        return doc

    # ── Format-Specific Parsers ────────────────────────────────────

    def _parse_pdf(self, file_path: str) -> ParsedDocument:
        """Parse PDF using pdfplumber for high-quality text extraction."""
        import pdfplumber

        pages_text = []
        page_count = 0

        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""

                # Also extract tables and convert to text
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        table_text = self._table_to_text(table)
                        if table_text and table_text not in text:
                            text += f"\n\n[Table]\n{table_text}\n"

                if text.strip():
                    pages_text.append(f"[Page {i + 1}]\n{text}")

        return ParsedDocument(
            content="\n\n".join(pages_text),
            metadata={"page_count": page_count},
            page_count=page_count,
        )

    def _parse_docx(self, file_path: str) -> ParsedDocument:
        """Parse DOCX files, extracting paragraphs, tables, and headers."""
        from docx import Document

        doc = Document(file_path)
        sections = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            # Preserve heading structure for context
            style = para.style.name if para.style else ""
            if "Heading" in style:
                level = style.replace("Heading ", "").strip()
                prefix = "#" * int(level) if level.isdigit() else "##"
                sections.append(f"{prefix} {text}")
            else:
                sections.append(text)

        # Extract tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                sections.append(f"[Table]\n" + "\n".join(rows))

        return ParsedDocument(
            content="\n\n".join(sections),
            metadata={"paragraph_count": len(doc.paragraphs)},
        )

    def _parse_pptx(self, file_path: str) -> ParsedDocument:
        """Parse PowerPoint presentations."""
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []

        for i, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_content.append(text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    slide_content.append(f"[Table]\n" + "\n".join(rows))

            if slide_content:
                slides_text.append(
                    f"[Slide {i + 1}]\n" + "\n".join(slide_content)
                )

        return ParsedDocument(
            content="\n\n".join(slides_text),
            metadata={"slide_count": len(prs.slides)},
            page_count=len(prs.slides),
        )

    def _parse_excel(self, file_path: str) -> ParsedDocument:
        """Parse Excel files — each sheet becomes a section."""
        import pandas as pd

        sheets_text = []
        xls = pd.ExcelFile(file_path)

        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            if df.empty:
                continue
            # Convert to readable text format
            text = f"[Sheet: {sheet_name}]\n"
            text += df.to_string(index=False, na_rep="")
            sheets_text.append(text)

        return ParsedDocument(
            content="\n\n".join(sheets_text),
            metadata={"sheet_count": len(xls.sheet_names)},
        )

    def _parse_csv(self, file_path: str) -> ParsedDocument:
        """Parse CSV files into readable text."""
        import pandas as pd

        df = pd.read_csv(file_path)
        content = df.to_string(index=False, na_rep="")

        return ParsedDocument(
            content=content,
            metadata={"row_count": len(df), "column_count": len(df.columns)},
        )

    def _parse_text(self, file_path: str) -> ParsedDocument:
        """Parse plain text files with encoding detection."""
        raw = open(file_path, "rb").read()
        detected = chardet.detect(raw)
        encoding = detected.get("encoding", "utf-8") or "utf-8"

        content = raw.decode(encoding, errors="replace")
        return ParsedDocument(content=content)

    def _parse_markdown(self, file_path: str) -> ParsedDocument:
        """Parse Markdown files — keep structure for better chunking."""
        raw = open(file_path, "rb").read()
        detected = chardet.detect(raw)
        encoding = detected.get("encoding", "utf-8") or "utf-8"

        content = raw.decode(encoding, errors="replace")
        return ParsedDocument(
            content=content,
            metadata={"format": "markdown"},
        )

    def _parse_html(self, file_path: str) -> ParsedDocument:
        """Parse HTML files — extract clean text from tags."""
        from bs4 import BeautifulSoup

        raw = open(file_path, "rb").read()
        detected = chardet.detect(raw)
        encoding = detected.get("encoding", "utf-8") or "utf-8"
        html = raw.decode(encoding, errors="replace")

        soup = BeautifulSoup(html, "html.parser")

        # Remove script and style elements
        for element in soup(["script", "style", "nav", "footer", "header"]):
            element.decompose()

        text = soup.get_text(separator="\n", strip=True)
        title = soup.title.string if soup.title else ""

        return ParsedDocument(
            content=text,
            metadata={"title": title},
        )

    def _parse_json(self, file_path: str) -> ParsedDocument:
        """Parse JSON files — flatten into readable key-value text."""
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        content = self._flatten_json(data)
        return ParsedDocument(content=content)

    # ── Helper Methods ─────────────────────────────────────────────

    def _table_to_text(self, table: list) -> str:
        """Convert a table (list of lists) to readable text."""
        if not table:
            return ""
        rows = []
        for row in table:
            if row:
                cells = [str(cell).strip() if cell else "" for cell in row]
                rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _flatten_json(self, data, prefix: str = "") -> str:
        """Recursively flatten JSON into readable text."""
        lines = []
        if isinstance(data, dict):
            for key, value in data.items():
                full_key = f"{prefix}.{key}" if prefix else key
                if isinstance(value, (dict, list)):
                    lines.append(self._flatten_json(value, full_key))
                else:
                    lines.append(f"{full_key}: {value}")
        elif isinstance(data, list):
            for i, item in enumerate(data):
                full_key = f"{prefix}[{i}]"
                if isinstance(item, (dict, list)):
                    lines.append(self._flatten_json(item, full_key))
                else:
                    lines.append(f"{full_key}: {item}")
        else:
            lines.append(f"{prefix}: {data}")

        return "\n".join(filter(None, lines))

    def _clean_text(self, text: str) -> str:
        """Clean extracted text: normalize whitespace, remove artifacts."""
        import re

        # Normalize whitespace
        text = re.sub(r"\t+", " ", text)
        # Remove excessive blank lines (keep max 2)
        text = re.sub(r"\n{4,}", "\n\n\n", text)
        # Remove null bytes and control characters (except newline/tab)
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
        # Trim
        text = text.strip()
        return text
